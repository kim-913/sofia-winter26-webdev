"""
Bubble Bliss — Clean Minimal Presentation (v4)
Rule: bullets ≤ 8 words · 18-20pt · generous boxes · audited before save
Run: python3 generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

# ── Palette ────────────────────────────────────────────────
PURPLE   = RGBColor(0x8B, 0x7B, 0xB5)
PURPLE_D = RGBColor(0x52, 0x46, 0x7A)
PURPLE_L = RGBColor(0xD4, 0xCE, 0xED)
PINK     = RGBColor(0xF4, 0xC7, 0xD6)
CREAM    = RGBColor(0xFA, 0xF8, 0xF5)
MINT     = RGBColor(0xB8, 0xE6, 0xD5)
MINT_D   = RGBColor(0x2A, 0x7A, 0x58)
PEACH    = RGBColor(0xFF, 0xCC, 0x99)
CHARCOAL = RGBColor(0x2C, 0x2C, 0x2C)
GRAY     = RGBColor(0x66, 0x66, 0x66)
LGRAY    = RGBColor(0xEC, 0xE9, 0xF5)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG  = RGBColor(0x1E, 0x1A, 0x2E)
CODE_GRN = RGBColor(0x9E, 0xE9, 0xC2)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Primitives ─────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, outline=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if outline: s.line.color.rgb = outline; s.line.width = Pt(0.5)
    else: s.line.fill.background()
    return s

def circle(slide, cx, cy, r, fill):
    s = slide.shapes.add_shape(9, cx-r, cy-r, r*2, r*2)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()

def txt(slide, text, l, t, w, h, sz=18, bold=False, italic=False,
        color=CHARCOAL, align=PP_ALIGN.LEFT):
    """Single text box. No auto-fit — keep text SHORT."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def bullets(slide, items, l, t, w, h, sz=18, color=CHARCOAL, gap=8):
    """Bullet list — items must be SHORT (≤ 8 words each)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = f"•  {item}"
        r.font.size = Pt(sz); r.font.color.rgb = color
    return tb

def code(slide, text, l, t, w, h, sz=11):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.name = "Courier New"
    r.font.color.rgb = CODE_GRN
    return tb

def pnum(slide, n, total=12):
    txt(slide, f"{n} / {total}", W-Inches(1.3), H-Inches(0.38),
        Inches(1.1), Inches(0.3), sz=11, color=GRAY, align=PP_ALIGN.RIGHT)

def topbar(slide, title, sub=None, tag=None, tag_color=None):
    rect(slide, 0, 0, W, Inches(0.055), PINK)
    rect(slide, 0, Inches(0.055), W, Inches(1.495), PURPLE)
    txt(slide, title, Inches(0.6), Inches(0.1), Inches(11.0), Inches(0.85),
        sz=34, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.6), Inches(0.94), Inches(11.0), Inches(0.45),
            sz=13, italic=True, color=RGBColor(0xDD, 0xD5, 0xFF))
    if tag:
        tc = tag_color or PINK
        p = slide.shapes.add_shape(1, Inches(0.6),
            Inches(0.88 if not sub else 0.88),
            Inches(len(tag)*0.095+0.25), Inches(0.26))
        p.fill.solid(); p.fill.fore_color.rgb = tc; p.line.fill.background()
        tf = p.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run(); r.text = tag
        r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = PURPLE_D if tc == PINK else WHITE
    rect(slide, 0, H-Inches(0.08), W, Inches(0.08), PINK)

def card(slide, l, t, w, h, accent=None):
    rect(slide, l, t, w, h, WHITE, outline=PURPLE_L)
    if accent:
        rect(slide, l, t, Inches(0.06), h, accent)

# ── Audit helper ───────────────────────────────────────────
def est_lines(text, box_w_in, sz_pt, gap_ratio=1.5):
    """Estimate rendered lines accounting for word-wrap."""
    chars_per_line = max(1, (box_w_in * 72) / (sz_pt * 0.55))
    lines = 0
    for para in text.split('\n'):
        if not para.strip():
            lines += 0.4
        else:
            lines += max(1.0, len(para) / chars_per_line)
    return lines

def audit(slide, slide_num):
    warnings = []
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        tf = shape.text_frame
        full = '\n'.join(p.text for p in tf.paragraphs).strip()
        if not full: continue
        try:
            sz = tf.paragraphs[0].runs[0].font.size
            sz_pt = sz.pt if sz else 14
        except: sz_pt = 14
        w_in = shape.width / 914400
        h_in = shape.height / 914400
        line_h = sz_pt / 72 * 1.55   # inches per line including gap
        est_h  = est_lines(full, w_in, sz_pt) * line_h
        if est_h > h_in + 0.15:
            warnings.append(
                f"  ⚠️  S{slide_num}: '{full[:55]}…'\n"
                f"     box={h_in:.2f}in, need≈{est_h:.2f}in (sz={sz_pt:.0f}pt, w={w_in:.1f}in)"
            )
    return warnings

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, CREAM)
rect(s1, 0, 0, Inches(8.8), H, PURPLE)
circle(s1, Inches(1.1), Inches(1.3), Inches(2.5), RGBColor(0x7A,0x6A,0x9F))
circle(s1, Inches(7.6), Inches(6.3), Inches(1.3), RGBColor(0x7A,0x6A,0x9F))
circle(s1, Inches(0.4), Inches(6.8), Inches(0.7), PINK)
circle(s1, Inches(11.0), Inches(1.8), Inches(1.5), LGRAY)
circle(s1, Inches(12.5), Inches(5.5), Inches(0.9), PEACH)
circle(s1, Inches(9.9),  Inches(6.4), Inches(0.45), MINT)

txt(s1, "Bubble",  Inches(0.7), Inches(0.6),  Inches(8.0), Inches(1.82), sz=82, bold=True, color=WHITE)
txt(s1, "Bliss",   Inches(0.7), Inches(2.28), Inches(8.0), Inches(1.82), sz=82, bold=True, color=PINK)
txt(s1, "Your Daily Dose of Happiness",
    Inches(0.7), Inches(3.4), Inches(7.8), Inches(0.5), sz=18, italic=True,
    color=RGBColor(0xE0,0xD8,0xFF))
rect(s1, Inches(0.7), Inches(4.05), Inches(4.0), Inches(0.04), PINK)
txt(s1, "Kim Ke", Inches(0.7), Inches(4.25), Inches(7.8), Inches(0.5),
    sz=22, bold=True, color=WHITE)
txt(s1, "Sofia University  ·  Web Development  ·  Winter 2026",
    Inches(0.7), Inches(4.82), Inches(7.8), Inches(0.4), sz=14,
    color=RGBColor(0xCC,0xC0,0xFF))
txt(s1, "March 22, 2026", Inches(0.7), Inches(5.28), Inches(4.0), Inches(0.38),
    sz=13, color=RGBColor(0xAA,0x98,0xDD))

# Right panel — just the key stats
txt(s1, "5 pages · 3 JS files · 0 frameworks",
    Inches(9.1), Inches(2.0), Inches(4.0), Inches(0.5), sz=15, bold=True, color=PURPLE_D)
txt(s1, "HTML5 · CSS3 · ES6+",
    Inches(9.1), Inches(2.6), Inches(4.0), Inches(0.45), sz=18, bold=True, color=CHARCOAL)
txt(s1, "Vanilla. No frameworks.",
    Inches(9.1), Inches(3.12), Inches(4.0), Inches(0.45), sz=14, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Design & Planning
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, CREAM)
topbar(s2, "Design & Planning", tag="Week 1 · 15 pts")

# Left col: audience + principles
card(s2, Inches(0.5), Inches(1.68), Inches(4.1), Inches(5.32), accent=PURPLE)
txt(s2, "Target Audience", Inches(0.72), Inches(1.82), Inches(3.7), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
bullets(s2,
    ["Students & young professionals, 18–30",
     "Boba enthusiasts in Palo Alto & Costa Mesa",
     "Mobile-first — ordering on the go",
     "Keyboard + screen reader users"],
    Inches(0.72), Inches(2.3), Inches(3.7), Inches(1.9),
    sz=15, gap=7)
txt(s2, "Design Principles", Inches(0.72), Inches(4.38), Inches(3.7), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
bullets(s2,
    ["Repetition — shared header/footer/cards",
     "Contrast — purple on cream, 4.5:1+",
     "Alignment — left-anchored, grid-based",
     "Proximity — related items in card groups"],
    Inches(0.72), Inches(4.85), Inches(3.7), Inches(1.9),
    sz=15, gap=7)

# Middle col: site map
card(s2, Inches(4.78), Inches(1.68), Inches(3.95), Inches(5.32), accent=MINT_D)
txt(s2, "Information Architecture",
    Inches(4.98), Inches(1.82), Inches(3.55), Inches(0.4), sz=15, bold=True, color=MINT_D)
# Site map visual
rect(s2, Inches(5.7), Inches(2.45), Inches(1.6), Inches(0.38), PURPLE)
txt(s2, "Home", Inches(5.7), Inches(2.47), Inches(1.6), Inches(0.34),
    sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s2, Inches(6.49), Inches(2.83), Inches(0.03), Inches(0.32), PURPLE_L)
rect(s2, Inches(4.98), Inches(3.15), Inches(3.5), Inches(0.03), PURPLE_L)
for j, (name, xoff) in enumerate([("Menu",Inches(4.98)),("About",Inches(5.73)),
                                    ("Locs",Inches(6.48)),("Contact",Inches(7.12))]):
    rect(s2, xoff+Inches(0.32), Inches(3.15), Inches(0.03), Inches(0.28), PURPLE_L)
    rect(s2, xoff, Inches(3.43), Inches(0.93), Inches(0.35), PURPLE_L, outline=PURPLE)
    txt(s2, name, xoff+Inches(0.03), Inches(3.45), Inches(0.87), Inches(0.3),
        sz=10, color=PURPLE_D, align=PP_ALIGN.CENTER)
txt(s2, "Flat hierarchy — all pages from nav",
    Inches(4.98), Inches(4.02), Inches(3.55), Inches(0.35), sz=12, italic=True, color=GRAY)
txt(s2, "Navigation", Inches(4.98), Inches(4.55), Inches(3.55), Inches(0.35),
    sz=15, bold=True, color=MINT_D)
bullets(s2,
    ["Semantic <nav> with aria-label",
     "Active page: aria-current='page'",
     "Skip-to-content link (every page)",
     "Hamburger: aria-expanded on toggle"],
    Inches(4.98), Inches(4.98), Inches(3.55), Inches(1.9),
    sz=14, gap=6)

# Right col: design system swatches
card(s2, Inches(8.91), Inches(1.68), Inches(3.95), Inches(5.32), accent=PEACH)
txt(s2, "Design System", Inches(9.11), Inches(1.82), Inches(3.55), Inches(0.4),
    sz=15, bold=True, color=RGBColor(0xB0,0x68,0x00))
txt(s2, "CSS Custom Properties", Inches(9.11), Inches(2.3), Inches(3.55), Inches(0.32),
    sz=12, bold=True, color=CHARCOAL)
swatches = [
    (PURPLE,   "#8B7BB5  Brand primary"),
    (PINK,     "#F4C7D6  Accents"),
    (MINT,     "#B8E6D5  Success"),
    (PEACH,    "#FFCC99  Badges"),
    (CHARCOAL, "#2C2C2C  Body text"),
]
for i, (clr, lbl) in enumerate(swatches):
    y = Inches(2.72 + i * 0.42)
    rect(s2, Inches(9.11), y, Inches(0.42), Inches(0.3), clr,
         outline=RGBColor(0xCC,0xC8,0xDC))
    txt(s2, lbl, Inches(9.61), y, Inches(3.0), Inches(0.3), sz=11.5, color=CHARCOAL)
txt(s2, "Typography", Inches(9.11), Inches(4.88), Inches(3.55), Inches(0.32),
    sz=12, bold=True, color=CHARCOAL)
txt(s2, "Georgia serif  /  Arial sans",
    Inches(9.11), Inches(5.28), Inches(3.55), Inches(0.32), sz=12, color=GRAY)
txt(s2, "Spacing: 8 · 16 · 24 · 32 · 48 · 64 px",
    Inches(9.11), Inches(5.65), Inches(3.55), Inches(0.32), sz=12, color=GRAY)
txt(s2, "Breakpoints: 480 · 768 · 1024 px",
    Inches(9.11), Inches(6.02), Inches(3.55), Inches(0.32), sz=12, color=GRAY)

pnum(s2, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Layout & Responsiveness
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, CREAM)
topbar(s3, "Layout & Responsiveness", tag="Week 2 · 25 pts")

concepts = [
    (PURPLE, "Semantic HTML",
     ["header · nav · main · section · article · footer",
      "table: caption + scope='col/row'",
      "form: fieldset + legend + label"]),
    (MINT_D, "CSS Box Model",
     ["box-sizing: border-box — global reset",
      "Spacing via tokens (--space-1 → --space-6)",
      "No magic pixel values"]),
    (PURPLE, "Flexbox",
     ["Nav links · hero · form buttons · footer",
      "Location cards side-by-side",
      "Vertical centering throughout"]),
    (MINT_D, "CSS Grid",
     ["Menu: auto-fill minmax(280px,1fr) → 1→3 col",
      "Values grid · team grid · footer 3-col",
      "Reflows automatically at breakpoints"]),
    (PURPLE, "Responsive (Mobile-First)",
     ["Base styles: 375px",
      "Breakpoints: 480px · 768px · 1024px",
      "@media (min-width: …) progressive layers"]),
    (MINT_D, "Advanced CSS",
     ["position: sticky — header stays at top",
      "tr:nth-child(even) — table zebra rows",
      "print.css — hides nav/hero for printing"]),
]

for i, (accent, title, pts) in enumerate(concepts):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.42)
    y = Inches(1.68 + row * 1.88)
    card(s3, x, y, Inches(6.15), Inches(1.74), accent=accent)
    txt(s3, title, x+Inches(0.2), y+Inches(0.1), Inches(5.75), Inches(0.38),
        sz=15, bold=True, color=CHARCOAL)
    bullets(s3, pts, x+Inches(0.2), y+Inches(0.56), Inches(5.75), Inches(1.1),
            sz=13.5, gap=4)

pnum(s3, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Interactivity: Menu Page
# ══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, CREAM)
topbar(s4, "Interactivity — Menu Page  🧋",
       sub="20 drinks · search · filter · modal · cart · favorites")

# Left bullets
card(s4, Inches(0.5), Inches(1.68), Inches(7.5), Inches(5.32), accent=PURPLE)
txt(s4, "What it does", Inches(0.7), Inches(1.82), Inches(7.1), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
bullets(s4,
    ["20 drinks stored as JS array of objects",
     "renderCards() builds all DOM dynamically",
     "Filter tabs — ARIA tablist/tab pattern",
     "Live search — 300ms debounce",
     "Modal — focus trap · Escape · aria-labelledby",
     "Cart — sessionStorage persists across pages",
     "Favorites — localStorage · greets you by name",
     "Listen — Web Speech API reads descriptions"],
    Inches(0.7), Inches(2.3), Inches(7.1), Inches(4.4),
    sz=16, gap=7)

# Right code
rect(s4, Inches(8.2), Inches(1.68), Inches(4.65), Inches(3.1), CODE_BG)
txt(s4, "// 300ms debounce", Inches(8.38), Inches(1.8), Inches(4.3), Inches(0.3),
    sz=11, bold=True, color=MINT)
code(s4,
"""let timer;
searchInput.addEventListener('input',
 () => {
   clearTimeout(timer);
   timer = setTimeout(() => {
     searchQuery = input.value
       .trim().toLowerCase();
     applyFilters();
   }, 300);
 }
);""",
    Inches(8.38), Inches(2.15), Inches(4.3), Inches(3.0), sz=12)

card(s4, Inches(8.2), Inches(4.92), Inches(4.65), Inches(2.08), accent=MINT_D)
txt(s4, "Progressive Enhancement",
    Inches(8.4), Inches(5.06), Inches(4.25), Inches(0.38), sz=14, bold=True, color=MINT_D)
bullets(s4,
    ["loading='lazy' on all images",
     "Search shows all drinks if JS fails",
     "Form submits without JS (server validates)"],
    Inches(8.4), Inches(5.52), Inches(4.25), Inches(1.35),
    sz=14, gap=7)

pnum(s4, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Accessibility
# ══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, CREAM)
topbar(s5, "Accessibility — WCAG POUR", tag="Week 3 · 30 pts")

wcag = [
    (PURPLE, "Perceivable",
     ["Alt text on all images",
      "Color contrast ≥ 4.5:1 for body text",
      "aria-live on cart count, errors, results"]),
    (MINT_D, "Operable",
     ["Skip-to-content link on every page",
      "Full keyboard navigation (Tab/Enter/Esc)",
      "Focus trap inside modal"]),
    (RGBColor(0xB0,0x68,0x00), "Understandable",
     ["lang='en' on <html>",
      "aria-current='page' on active nav link",
      "aria-describedby links inputs to errors"]),
    (PURPLE_D, "Robust",
     ["Semantic <button> — never <div onclick>",
      "role=dialog · aria-modal on overlay",
      "role=tablist/tab on filter buttons"]),
]

for i, (accent, title, pts) in enumerate(wcag):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.42)
    y = Inches(1.68 + row * 2.68)
    card(s5, x, y, Inches(6.15), Inches(2.52), accent=accent)
    txt(s5, title, x+Inches(0.2), y+Inches(0.12), Inches(5.75), Inches(0.42),
        sz=18, bold=True, color=CHARCOAL)
    bullets(s5, pts, x+Inches(0.2), y+Inches(0.64), Inches(5.75), Inches(1.75),
            sz=15, gap=8)

pnum(s5, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Contact Form
# ══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, CREAM)
topbar(s6, "Contact Form & Validation", tag="Week 3 · 30 pts")

# Left: input types
card(s6, Inches(0.5), Inches(1.68), Inches(5.6), Inches(5.32), accent=PURPLE)
txt(s6, "All Input Types", Inches(0.7), Inches(1.82), Inches(5.2), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
bullets(s6,
    ["text — full name",
     "email — address with regex check",
     "tel — phone (optional)",
     "radio — inquiry type (pill selector)",
     "checkbox — 'how did you find us?'",
     "select — preferred location",
     "textarea — message, 500-char limit",
     "fieldset + legend on both groups"],
    Inches(0.7), Inches(2.3), Inches(5.2), Inches(4.5),
    sz=15, gap=6)

# Right: validation + submit flow
card(s6, Inches(6.3), Inches(1.68), Inches(6.55), Inches(2.4), accent=MINT_D)
txt(s6, "Live Validation", Inches(6.5), Inches(1.82), Inches(6.15), Inches(0.4),
    sz=15, bold=True, color=MINT_D)
bullets(s6,
    ["novalidate — JS handles all validation",
     "blur event: validate when user leaves field",
     "input event: clear error as user types",
     "aria-describedby links input → error span"],
    Inches(6.5), Inches(2.3), Inches(6.15), Inches(1.6),
    sz=14, gap=6)

card(s6, Inches(6.3), Inches(4.22), Inches(6.55), Inches(2.78), accent=PURPLE)
txt(s6, "Submit Experience", Inches(6.5), Inches(4.36), Inches(6.15), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
# Flow steps as stacked boxes
steps = [
    ("Click 'Send Message'",             LGRAY,    CHARCOAL),
    ("JS validates all fields",          LGRAY,    CHARCOAL),
    ("❌  Shake button + scroll to error", RGBColor(0xF8,0xD7,0xD7), RGBColor(0x8B,0x00,0x00)),
    ("✅  Spinner: 'Sending…' (1.4 s)",   RGBColor(0xD4,0xED,0xDA), MINT_D),
    ("Typewriter: 'Hey Kim! …'",         PURPLE_L, PURPLE_D),
]
for j, (label, bg_c, fc) in enumerate(steps):
    rect(s6, Inches(6.5), Inches(4.9+j*0.38), Inches(6.0), Inches(0.34), bg_c)
    txt(s6, label, Inches(6.6), Inches(4.92+j*0.38), Inches(5.8), Inches(0.3),
        sz=12, color=fc)

pnum(s6, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Locations Table
# ══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7, CREAM)
topbar(s7, "Locations — Accessible Data Table",
       sub="Semantic table markup · CSS styling · JavaScript enhancement")

card(s7, Inches(0.5), Inches(1.68), Inches(5.6), Inches(5.32), accent=MINT_D)
txt(s7, "Table Structure", Inches(0.7), Inches(1.82), Inches(5.2), Inches(0.4),
    sz=15, bold=True, color=MINT_D)
bullets(s7,
    ["<caption> names the table",
     "<th scope='col'> — Day, Downtown, Campus",
     "<th scope='row'> — each day of the week",
     "Screen readers: 'Monday, Downtown: 9AM–9PM'",
     "Zebra rows: tr:nth-child(even)",
     "Today's row: new Date().getDay() → pink",
     "Location cards: Flexbox side-by-side layout"],
    Inches(0.7), Inches(2.3), Inches(5.2), Inches(4.5),
    sz=15, gap=7)

# Table visual mock
card(s7, Inches(6.4), Inches(1.68), Inches(6.45), Inches(5.32), accent=PURPLE)
txt(s7, "Preview", Inches(6.6), Inches(1.82), Inches(6.05), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
rect(s7, Inches(6.6), Inches(2.35), Inches(6.0), Inches(0.38), PURPLE)
for j, h in enumerate(["Day", "Palo Alto", "Costa Mesa"]):
    txt(s7, h, Inches(6.63+j*2.0), Inches(2.37), Inches(1.85), Inches(0.34),
        sz=12, bold=True, color=WHITE)
rows_data = [
    ("Mon – Fri",  "9 AM – 9 PM",  "9 AM – 9 PM",  WHITE),
    ("Saturday",   "10 AM – 10 PM","10 AM – 10 PM", LGRAY),
    ("Sunday",     "10 AM – 10 PM","10 AM – 10 PM", WHITE),
    ("Today ★",   "9 AM – 9 PM",  "9 AM – 9 PM",   PINK),
]
for ri, (d, pa, cm, rbg) in enumerate(rows_data):
    ry = Inches(2.73 + ri*0.43)
    rect(s7, Inches(6.6), ry, Inches(6.0), Inches(0.4), rbg, outline=PURPLE_L)
    for ci, cell in enumerate([d, pa, cm]):
        txt(s7, cell, Inches(6.63+ci*2.0), ry+Inches(0.03), Inches(1.85), Inches(0.34),
            sz=12, bold=(ci==0), color=PURPLE_D if ri==3 else CHARCOAL)
txt(s7, "★  highlighted via new Date().getDay()",
    Inches(6.6), Inches(4.48), Inches(6.0), Inches(0.32), sz=11, italic=True, color=GRAY)

pnum(s7, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Apple Scroll Effects
# ══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
bg(s8, CREAM)
topbar(s8, "Scroll Experience — Beyond the Rubric",
       sub="Pure CSS + vanilla JS · no libraries")

effects = [
    (PURPLE,                    "Clip-Reveal Headings",
     ["overflow:hidden clips the container",
      "Inner span: translateY(105%) → hidden",
      "IntersectionObserver → snap to 0"]),
    (MINT_D,                    "Word-by-Word Title",
     ["JS splits h1 into one span per word",
      "--i CSS var sets animation-delay per word",
      "Each word rises in sequence"]),
    (RGBColor(0xB0,0x68,0x00), "Hero Parallax",
     ["Scroll event moves content at 22% speed",
      "Opacity fades as you scroll away",
      "8 lines of JS — zero libraries"]),
    (PURPLE_D,                  "3D Card Tilt",
     ["mousemove → perspective(700px) rotateX/Y",
      "Event delegation — works on dynamic cards",
      "Spring ease on reset"]),
    (MINT_D,                    "Staggered Entrance",
     ["--stagger-i injected on grid children",
      "transition-delay: calc(i * 0.12s)",
      "Wave cascade effect on scroll"]),
    (PURPLE,                    "Count-Up Stats",
     ["IntersectionObserver triggers on stat",
      "requestAnimationFrame loop 0→target",
      "Cubic ease-out for natural deceleration"]),
]

for i, (accent, title, pts) in enumerate(effects):
    col = i % 2; row = i // 2
    x = Inches(0.5 + col * 6.42)
    y = Inches(1.68 + row * 1.88)
    card(s8, x, y, Inches(6.15), Inches(1.74), accent=accent)
    txt(s8, title, x+Inches(0.2), y+Inches(0.1), Inches(5.75), Inches(0.38),
        sz=15, bold=True, color=CHARCOAL)
    bullets(s8, pts, x+Inches(0.2), y+Inches(0.55), Inches(5.75), Inches(1.1),
            sz=13.5, gap=4)

pnum(s8, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Key Design Decisions
# ══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
bg(s9, CREAM)
topbar(s9, "Key Design Decisions",
       sub="The reasoning behind the architecture")

decisions = [
    (PURPLE,                    "No frameworks",
     ["Shows real DOM/event/state mastery",
      "Frameworks hide what the rubric tests"]),
    (MINT_D,                    "Design system first",
     ["All 5 pages consistent automatically",
      "One token change updates everything"]),
    (RGBColor(0xB0,0x68,0x00), "sessionStorage vs localStorage",
     ["Cart clears on session end (correct)",
      "Favorites persist — personalization"]),
    (PURPLE_D,                  "Data array drives the menu",
     ["renderCards() = pure function",
      "Add a drink = one JS object"]),
    (PURPLE,                    "Progressive enhancement",
     ["Form works without JavaScript",
      "JS layers on top — not required"]),
    (MINT_D,                    "ARIA designed in",
     ["Not retrofitted — built into structure",
      "JS manages dynamic ARIA states"]),
]

for i, (accent, title, pts) in enumerate(decisions):
    col = i % 2; row = i // 3
    x = Inches(0.5 + col * 6.42)
    y = Inches(1.68 + row * 1.88)
    card(s9, x, y, Inches(6.15), Inches(1.74), accent=accent)
    txt(s9, title, x+Inches(0.2), y+Inches(0.1), Inches(5.75), Inches(0.38),
        sz=16, bold=True, color=CHARCOAL)
    bullets(s9, pts, x+Inches(0.2), y+Inches(0.57), Inches(5.75), Inches(1.05),
            sz=15, gap=7)

pnum(s9, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Performance & Documentation
# ══════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
bg(s10, CREAM)
topbar(s10, "Performance & Documentation", tag="Week 4 · 20 pts")

card(s10, Inches(0.5), Inches(1.68), Inches(5.95), Inches(5.32), accent=PURPLE)
txt(s10, "Performance", Inches(0.7), Inches(1.82), Inches(5.55), Inches(0.4),
    sz=15, bold=True, color=PURPLE_D)
bullets(s10,
    ["loading='lazy' on every image",
     "Unsplash: ?w=600&q=80 (server resize)",
     "{ passive: true } on scroll/mousemove",
     "CSS transform/opacity only — no reflow",
     "unobserve() after first trigger",
     "All JS loaded at end of <body>"],
    Inches(0.7), Inches(2.3), Inches(5.55), Inches(4.4),
    sz=15, gap=8)

card(s10, Inches(6.65), Inches(1.68), Inches(6.2), Inches(2.55), accent=MINT_D)
txt(s10, "README.md", Inches(6.85), Inches(1.82), Inches(5.8), Inches(0.4),
    sz=15, bold=True, color=MINT_D)
bullets(s10,
    ["Setup: open .html in browser — zero build",
     "Features per page",
     "Accessibility considerations",
     "Known limitations"],
    Inches(6.85), Inches(2.3), Inches(5.8), Inches(1.75),
    sz=15, gap=7)

card(s10, Inches(6.65), Inches(4.38), Inches(6.2), Inches(2.62), accent=PEACH)
txt(s10, "Design Dossier", Inches(6.85), Inches(4.52), Inches(5.8), Inches(0.4),
    sz=15, bold=True, color=RGBColor(0xB0,0x68,0x00))
bullets(s10,
    ["Audience + persona",
     "Site map + wireframes (desktop & mobile)",
     "Color palette with contrast rationale",
     "Accessibility checklist + fixes"],
    Inches(6.85), Inches(5.0), Inches(5.8), Inches(1.78),
    sz=15, gap=7)

pnum(s10, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Challenges & Lessons
# ══════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
bg(s11, CREAM)
topbar(s11, "Challenges & What I Learned")

items = [
    (PURPLE,                    "Design system before pages",
     "Forcing every value through CSS tokens felt slow at first — then the 5th page took 20% of the time the 1st did."),
    (MINT_D,                    "Accessibility is structure",
     "Focus traps, scope attributes, ARIA live regions can't be added at the end. They have to be in the HTML from the start."),
    (RGBColor(0xB0,0x68,0x00), "Real bugs teach real concepts",
     "A top-level return in strict mode silently kills the whole script. Debugging that taught me more than any tutorial."),
    (PURPLE_D,                  "Event delegation handles the future",
     "Attaching listeners to document.closest('.card') works for cards that don't exist yet — because menu.js generates them dynamically."),
]

for i, (accent, title, detail) in enumerate(items):
    col = i % 2; row = i // 2
    x = Inches(0.5 + col * 6.42)
    y = Inches(1.68 + row * 2.68)
    card(s11, x, y, Inches(6.15), Inches(2.52), accent=accent)
    txt(s11, title, x+Inches(0.2), y+Inches(0.12), Inches(5.75), Inches(0.42),
        sz=16, bold=True, color=CHARCOAL)
    txt(s11, detail, x+Inches(0.2), y+Inches(0.66), Inches(5.75), Inches(1.65),
        sz=13, color=GRAY)

pnum(s11, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Q&A
# ══════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
bg(s12, PURPLE)
circle(s12, Inches(1.6),  Inches(5.5), Inches(3.0), RGBColor(0x7A,0x6A,0xA2))
circle(s12, Inches(12.2), Inches(1.0), Inches(1.9), RGBColor(0x7A,0x6A,0xA2))
circle(s12, Inches(13.2), Inches(7.1), Inches(1.3), PINK)
circle(s12, Inches(0.5),  Inches(0.8), Inches(0.6), PEACH)

txt(s12, "Let's see it live.",
    Inches(1.0), Inches(1.1), Inches(11.3), Inches(1.8),
    sz=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s12, "Questions welcome  🧋",
    Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.75),
    sz=28, bold=True, color=PINK, align=PP_ALIGN.CENTER)

rect(s12, Inches(1.5), Inches(4.1), Inches(10.3), Inches(1.85),
     RGBColor(0x62,0x52,0x8A))
txt(s12, "Demo order",
    Inches(1.72), Inches(4.22), Inches(2.0), Inches(0.38),
    sz=13, bold=True, color=PINK)
bullets(s12,
    ["Home → scroll to see clip-reveal + parallax",
     "Menu → search · filter · open card · add to cart",
     "Contact → fill form · submit → typewriter",
     "Locations → today's row highlighted"],
    Inches(1.72), Inches(4.65), Inches(9.9), Inches(1.2),
    sz=14, color=WHITE, gap=3)

txt(s12, "Kim Ke  ·  Sofia University  ·  Web Development  ·  Winter 2026",
    Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.45),
    sz=13, color=RGBColor(0xCC,0xC0,0xFF), align=PP_ALIGN.CENTER)
rect(s12, 0, H-Inches(0.08), W, Inches(0.08), PINK)

# ── Audit before save ──────────────────────────────────────
print("\n── Overflow audit ──────────────────────────────────")
all_warnings = []
for i, slide in enumerate(prs.slides):
    all_warnings.extend(audit(slide, i+1))
if all_warnings:
    for w in all_warnings:
        print(w)
    print(f"\n{len(all_warnings)} potential overflow(s) found.")
else:
    print("✅  No overflow detected.")

# ── Save ───────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Bubble_Bliss_Presentation.pptx")
prs.save(out)
print(f"\n✅  Saved → {out}")
